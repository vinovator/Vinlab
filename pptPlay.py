# pptPlay.py
# Python 2.76

# Importing python-pptx module
from pptx import Presentation

# This is the template based on which PPT will be created
# If None is passed then blank ppt with no slides will be created
PPT_TEMPLATE = None

# The updated file will be saved in this name
TARGET_PPT = "Python_PPT.pptx"


def main():
	# Create a ppt based on template
	prs = Presentation(PPT_TEMPLATE)
	
	# Iterating through each layout present in the file
	# Different PPT files have different number of layouts
	# "__" is used for unnecessary variables that are not used anywhere
	for layout_index, __ in enumerate(prs.slide_layouts):
		# define a slide layout
		slide_layout = prs.slide_layouts[layout_index]
		
		# Assign the layout to a slide
		slide = prs.slides.add_slide(slide_layout)
		
		# Once a slide is created, everything else is done through "shapes"
		# Shapes behave the same way as a list
		
		# Assign the layout name as the slide title
		# Some slide layouts do not have slide title shape
		if(slide.shapes.title is not None):
			slide.shapes.title.text = slide_layout.name
		
		# Find out how many objects are present in the slide to playwith
		print("There are {0} shapes and {1} placeholders in the slide layout {2}".format(len(slide.shapes), len(slide.placeholders), slide_layout.name))
		
		# Iterate over the placeholders in the slide
		# "__" is used for unnecessary variables that are not used anywhere
		for placeholder_index, __ in enumerate(slide.shapes):
			# If the placeholder is not the slide title, then update it with placeholder type
			if(slide.shapes[placeholder_index].has_text_frame):
				# The name for title placeholder is always "Title 1", title placeholder is already populated with layout name; so skip
				if not(slide.shapes[placeholder_index].name in "Title 1"):
					slide.shapes[placeholder_index].text = slide.shapes[placeholder_index].name
	
	
	# Finally save the presentation
	prs.save(TARGET_PPT)
	
	print("Done. {0} slide(s) added".format(len(prs.slides)))

if __name__ == "__main__":
	main()